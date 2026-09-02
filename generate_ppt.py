"""
generate_ppt.py
Deepfake Investigation System — Professional PowerPoint Generator
Produces a fully styled, multi-slide presentation using python-pptx.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree

# ─── Color Palette ────────────────────────────────────────────────────────────
C_BG_DARK     = RGBColor(0x0A, 0x0F, 0x1D)   # #0a0f1d
C_BG_CARD     = RGBColor(0x11, 0x18, 0x27)   # #111827
C_ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # #3b82f6
C_ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)   # #06b6d4
C_ACCENT_EMR  = RGBColor(0x10, 0xB9, 0x81)   # #10b981 emerald
C_ACCENT_ROSE = RGBColor(0xF4, 0x3F, 0x5E)   # #f43f5e rose
C_ACCENT_AMB  = RGBColor(0xF5, 0x9E, 0x0B)   # #f59e0b amber
C_TEXT_MAIN   = RGBColor(0xF8, 0xFA, 0xFC)   # #f8fafc
C_TEXT_MUTED  = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_DIVIDER     = RGBColor(0x1E, 0x29, 0x3B)   # subtle card border

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ARCH_IMG = os.path.join(BASE_DIR, "architecture.png")
OUTPUT   = os.path.join(BASE_DIR, "Deepfake_Investigation_System.pptx")


# ─── Helper: fill slide background ────────────────────────────────────────────
def fill_bg(slide, color: RGBColor):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─── Helper: add rectangle shape ──────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()          # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


# ─── Helper: add text box ─────────────────────────────────────────────────────
def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_TEXT_MAIN,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


# ─── Helper: accent line ──────────────────────────────────────────────────────
def add_line(slide, left, top, width, color: RGBColor, thickness=Pt(2)):
    line = slide.shapes.add_shape(1, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


# ─── Helper: bullet list ──────────────────────────────────────────────────────
def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=C_TEXT_MUTED, bullet_color=C_ACCENT_CYAN,
                title=None, title_color=C_TEXT_MAIN, title_size=20):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    start = 0
    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"
        start = 1

    for i, item in enumerate(items):
        if i == 0 and not title:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    layout = prs.slide_layouts[6]          # blank
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    # Gradient-style top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_BLUE)

    # Left accent panel
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_BG_CARD)

    # Left vivid accent strip
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_CYAN)

    # Shield / emoji icon area
    add_text(slide, "Shield", Inches(1.2), Inches(1.3), Inches(2), Inches(1.5),
             font_size=60, align=PP_ALIGN.CENTER)

    # Main title
    add_text(slide, "DEEPFAKE INVESTIGATION SYSTEM",
             Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.1),
             font_size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Accent line under title
    add_line(slide, Inches(3.0), Inches(3.75), Inches(7.3), C_ACCENT_CYAN, Pt(3))

    # Subtitle
    add_text(slide,
             "Advanced AI-Powered Digital Forensics & Media Authenticity Verification Platform",
             Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.7),
             font_size=17, italic=True, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Tech tags row
    tags = [
        ("EfficientNet-B3", C_ACCENT_BLUE),
        ("Grad-CAM XAI", C_ACCENT_CYAN),
        ("OpenCV", C_ACCENT_EMR),
        ("Flask REST API", C_ACCENT_AMB),
        ("ReportLab PDF", C_ACCENT_ROSE),
    ]
    tag_x = Inches(1.3)
    for label, col in tags:
        box = add_rect(slide, tag_x, Inches(4.85), Inches(1.9), Inches(0.42), col)
        add_text(slide, label, tag_x, Inches(4.87), Inches(1.9), Inches(0.42),
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tag_x += Inches(2.05)

    # Bottom footer
    add_text(slide, "Digital Forensics  |  AI Research  |  Media Verification",
             Inches(0), Inches(7.05), SLIDE_W, Inches(0.4),
             font_size=11, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 0, Inches(7.44), SLIDE_W, Inches(0.06), C_ACCENT_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
def slide_problem(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_ROSE)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_ROSE)

    # Slide number badge
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_ROSE)
    add_text(slide, "02", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "THE PROBLEM", Inches(0.4), Inches(0.2), Inches(6), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_ROSE, align=PP_ALIGN.LEFT)
    add_text(slide, "Deepfakes: A Growing Threat to Digital Trust",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN, align=PP_ALIGN.LEFT)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_ROSE, Pt(2))

    # Stat cards
    stats = [
        ("900%", "Increase in deepfake content\nonline since 2019", C_ACCENT_ROSE),
        ("96%", "Of deepfakes target\nindividuals non-consensually", C_ACCENT_AMB),
        ("$25B+", "Projected financial fraud\ndamage by 2027", C_ACCENT_CYAN),
        ("3.8s", "Average time to generate\na convincing deepfake", C_ACCENT_EMR),
    ]
    card_x = Inches(0.4)
    for val, desc, col in stats:
        add_rect(slide, card_x, Inches(1.6), Inches(2.95), Inches(1.65), C_BG_CARD)
        add_line(slide, card_x, Inches(1.6), Inches(2.95), col, Pt(4))
        add_text(slide, val, card_x, Inches(1.7), Inches(2.95), Inches(0.75),
                 font_size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, desc, card_x, Inches(2.42), Inches(2.95), Inches(0.8),
                 font_size=13, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
        card_x += Inches(3.12)

    # Challenges
    add_text(slide, "Key Challenges", Inches(0.4), Inches(3.5), Inches(6), Inches(0.45),
             font_size=20, bold=True, color=C_TEXT_MAIN)
    challenges = [
        "Hyper-realistic AI face-swaps that fool the human eye",
        "No reliable manual verification method at scale",
        "Criminal use in disinformation campaigns & identity fraud",
        "Evidence integrity compromised in legal proceedings",
        "Rapid evolution of generative AI (GANs, Diffusion Models)",
    ]
    add_bullets(slide, challenges, Inches(0.4), Inches(3.95), Inches(6.2), Inches(2.8),
                font_size=15, color=C_TEXT_MUTED)

    # Solution teaser box
    add_rect(slide, Inches(6.9), Inches(3.45), Inches(6.0), Inches(3.7), C_BG_CARD)
    add_line(slide, Inches(6.9), Inches(3.45), Inches(6.0), C_ACCENT_ROSE, Pt(3))
    add_text(slide, "Our Solution",
             Inches(7.1), Inches(3.6), Inches(5.6), Inches(0.45),
             font_size=18, bold=True, color=C_ACCENT_CYAN)
    add_text(slide,
             "An enterprise-grade AI forensic platform that automatically\n"
             "detects synthetic facial manipulation with:\n\n"
             "- Deep CNN (EfficientNet-B3) binary classification\n"
             "- Explainable AI heatmaps via Grad-CAM\n"
             "- Frame-level video forensics\n"
             "- Cryptographically-structured forensic reports\n"
             "- Real-time REST API integration",
             Inches(7.1), Inches(4.1), Inches(5.6), Inches(2.7),
             font_size=14, color=C_TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════════
def slide_objectives(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_CYAN)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_CYAN)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_CYAN)
    add_text(slide, "03", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "PROJECT OBJECTIVES", Inches(0.4), Inches(0.2), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_CYAN)
    add_text(slide, "What We Set Out to Achieve",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_CYAN, Pt(2))

    objectives = [
        ("Target", "Automated Detection",
         "Build an end-to-end pipeline that automatically identifies synthetic facial manipulation in uploaded images and videos without manual intervention.",
         C_ACCENT_BLUE),
        ("Lens", "Explainable AI",
         "Provide transparency into model decisions using Grad-CAM activation heatmaps that highlight the exact pixel regions triggering deepfake classification.",
         C_ACCENT_CYAN),
        ("Film", "Video Forensics",
         "Implement temporal frame-by-frame analysis for video files, computing per-frame confidence and aggregate fake-ratio metrics.",
         C_ACCENT_EMR),
        ("Doc", "Forensic Evidence",
         "Auto-generate structured PDF investigation reports suitable for legal and evidentiary use, including confidence scores, heatmaps and risk assessments.",
         C_ACCENT_AMB),
        ("Chart", "Analytics Platform",
         "Provide a real-time interactive dashboard with historical case tracking, trend visualisations, and 1-click CSV export.",
         C_ACCENT_ROSE),
        ("API", "REST API",
         "Expose a headless JSON REST API endpoint enabling seamless integration with third-party microservices and mobile applications.",
         C_ACCENT_BLUE),
    ]

    cols = 3
    card_w = Inches(4.1)
    card_h = Inches(1.85)
    pad_x  = Inches(0.35)
    pad_y  = Inches(1.6)
    gap_x  = Inches(0.3)
    gap_y  = Inches(0.25)

    for i, (icon, title, desc, col) in enumerate(objectives):
        row = i // cols
        col_i = i % cols
        x = pad_x + col_i * (card_w + gap_x)
        y = pad_y + row * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h, C_BG_CARD)
        add_line(slide, x, y, card_w, col, Pt(4))
        add_text(slide, f"[{icon}]", x + Inches(0.12), y + Inches(0.08),
                 Inches(0.5), Inches(0.5), font_size=14, color=col)
        add_text(slide, title, x + Inches(0.65), y + Inches(0.1),
                 card_w - Inches(0.8), Inches(0.4),
                 font_size=15, bold=True, color=col)
        add_text(slide, desc, x + Inches(0.12), y + Inches(0.52),
                 card_w - Inches(0.2), card_h - Inches(0.65),
                 font_size=12, color=C_TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_architecture(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_BLUE)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_BLUE)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_BLUE)
    add_text(slide, "04", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "SYSTEM ARCHITECTURE", Inches(0.4), Inches(0.2), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_BLUE)
    add_text(slide, "7-Stage Detection Pipeline",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_BLUE, Pt(2))

    # Architecture image
    if os.path.exists(ARCH_IMG):
        slide.shapes.add_picture(ARCH_IMG, Inches(0.5), Inches(1.55), Inches(3.8), Inches(5.6))

    # Pipeline stages as cards
    stages = [
        ("01", "Input Validation & MIME Guard", "File type and size verification before processing.", C_ACCENT_BLUE),
        ("02", "OpenCV Face Extraction", "Haar Cascade face detection + proportional padding crop.", C_ACCENT_CYAN),
        ("03", "EfficientNet-B3 Inference", "300x300 input -> deep CNN binary classification.", C_ACCENT_EMR),
        ("04", "Grad-CAM XAI Heatmap", "Gradient-weighted activation maps on last conv layer.", C_ACCENT_AMB),
        ("05", "Risk & Confidence Scoring", "Tiered risk: Low -> Medium -> High -> Very High.", C_ACCENT_ROSE),
        ("06", "SQLite Case Logging", "Persistent investigation history with SQLAlchemy ORM.", C_ACCENT_BLUE),
        ("07", "PDF Report & Dashboard", "ReportLab PDF + Chart.js real-time analytics.", C_ACCENT_CYAN),
    ]

    s_w = Inches(4.65)
    s_h = Inches(0.72)
    s_x = Inches(4.55)
    s_y = Inches(1.55)
    gap = Inches(0.05)

    for num, title, desc, col in stages:
        add_rect(slide, s_x, s_y, s_w, s_h, C_BG_CARD)
        # Number badge
        add_rect(slide, s_x, s_y, Inches(0.5), s_h, col)
        add_text(slide, num, s_x, s_y, Inches(0.5), s_h,
                 font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, s_x + Inches(0.58), s_y + Inches(0.04),
                 s_w - Inches(0.65), Inches(0.35),
                 font_size=13, bold=True, color=C_TEXT_MAIN)
        add_text(slide, desc, s_x + Inches(0.58), s_y + Inches(0.38),
                 s_w - Inches(0.65), Inches(0.3),
                 font_size=11, color=C_TEXT_MUTED)
        s_y += s_h + gap

    # Right-side tech stack
    tech_x = Inches(9.45)
    tech_y = Inches(1.55)
    add_rect(slide, tech_x, tech_y, Inches(3.5), Inches(5.6), C_BG_CARD)
    add_line(slide, tech_x, tech_y, Inches(3.5), C_ACCENT_CYAN, Pt(3))
    add_text(slide, "Technology Stack", tech_x + Inches(0.15), tech_y + Inches(0.1),
             Inches(3.2), Inches(0.4), font_size=14, bold=True, color=C_ACCENT_CYAN)
    techs = [
        ("Python 3.11+",       "Runtime"),
        ("TensorFlow / Keras",  "Deep Learning"),
        ("EfficientNet-B3",    "CNN Backbone"),
        ("OpenCV 4",           "Computer Vision"),
        ("Flask 3",            "Web Framework"),
        ("SQLAlchemy / SQLite","Database ORM"),
        ("ReportLab",          "PDF Generation"),
        ("Chart.js",           "Analytics UI"),
        ("Docker",             "Containerization"),
        ("Render.com",         "Cloud Deployment"),
    ]
    t_y = tech_y + Inches(0.55)
    for tech, role in techs:
        add_text(slide, f">> {tech}", tech_x + Inches(0.15), t_y,
                 Inches(2.0), Inches(0.38), font_size=12, bold=True, color=C_TEXT_MAIN)
        add_text(slide, role, tech_x + Inches(2.2), t_y,
                 Inches(1.1), Inches(0.38), font_size=11, color=C_TEXT_MUTED, align=PP_ALIGN.RIGHT)
        t_y += Inches(0.47)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DEEP LEARNING MODEL
# ═══════════════════════════════════════════════════════════════════════════════
def slide_model(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_EMR)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_EMR)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_EMR)
    add_text(slide, "05", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "DEEP LEARNING MODEL", Inches(0.4), Inches(0.2), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_EMR)
    add_text(slide, "EfficientNet-B3 CNN Architecture",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_EMR, Pt(2))

    # Model info cards - left column
    model_props = [
        ("Input Shape", "300 x 300 x 3 (RGB)", C_ACCENT_BLUE),
        ("Architecture", "EfficientNet-B3 + Custom Head", C_ACCENT_CYAN),
        ("Task", "Binary Classification", C_ACCENT_EMR),
        ("Output", "Sigmoid -> Fake / Real", C_ACCENT_AMB),
        ("Preprocessing", "EfficientNet preprocess_input()", C_ACCENT_ROSE),
        ("Optimizer", "Adam + Binary Cross-Entropy", C_ACCENT_BLUE),
    ]

    p_y = Inches(1.6)
    for label, val, col in model_props:
        add_rect(slide, Inches(0.4), p_y, Inches(4.1), Inches(0.65), C_BG_CARD)
        add_line(slide, Inches(0.4), p_y, Inches(0.06), col, Pt(65))
        add_text(slide, label, Inches(0.65), p_y + Inches(0.03),
                 Inches(1.5), Inches(0.3), font_size=11, color=C_TEXT_MUTED)
        add_text(slide, val, Inches(0.65), p_y + Inches(0.32),
                 Inches(3.7), Inches(0.3), font_size=13, bold=True, color=C_TEXT_MAIN)
        p_y += Inches(0.7)

    # Face pipeline flow diagram
    add_rect(slide, Inches(4.8), Inches(1.6), Inches(8.15), Inches(5.65), C_BG_CARD)
    add_line(slide, Inches(4.8), Inches(1.6), Inches(8.15), C_ACCENT_EMR, Pt(3))
    add_text(slide, "Inference Pipeline",
             Inches(4.95), Inches(1.7), Inches(7), Inches(0.4),
             font_size=15, bold=True, color=C_ACCENT_EMR)

    pipeline_steps = [
        ("Input Image", "Raw JPEG/PNG upload", C_ACCENT_BLUE),
        ("Face Detection", "Haar Cascade -> largest face ROI + 20% padding", C_ACCENT_CYAN),
        ("Resize", "cv2.resize -> 300x300  BGR->RGB conversion", C_ACCENT_EMR),
        ("Normalize", "EfficientNet preprocess_input() -> float32", C_ACCENT_AMB),
        ("Forward Pass", "model.predict(batch) -> score in [0, 1]", C_ACCENT_ROSE),
        ("Decision", "score < 0.45 = Fake | > 0.55 = Real | else Uncertain", C_ACCENT_CYAN),
        ("Risk Tier", "Confidence -> Low / Medium / High / Very High", C_ACCENT_EMR),
        ("Grad-CAM", "Gradient of class score wrt last conv layer activations", C_ACCENT_BLUE),
    ]

    step_y = Inches(2.2)
    for i, (name, detail, col) in enumerate(pipeline_steps):
        add_rect(slide, Inches(4.95), step_y, Inches(7.8), Inches(0.53), C_DIVIDER)
        add_rect(slide, Inches(4.95), step_y, Inches(0.07), Inches(0.53), col)
        add_text(slide, f"{i+1:02d}  {name}",
                 Inches(5.12), step_y + Inches(0.01),
                 Inches(1.9), Inches(0.27), font_size=12, bold=True, color=col)
        add_text(slide, detail,
                 Inches(7.1), step_y + Inches(0.01),
                 Inches(5.5), Inches(0.5), font_size=11, color=C_TEXT_MUTED)
        step_y += Inches(0.56)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ═══════════════════════════════════════════════════════════════════════════════
def slide_features(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_AMB)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_AMB)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_AMB)
    add_text(slide, "06", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "KEY FEATURES", Inches(0.4), Inches(0.2), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_AMB)
    add_text(slide, "Platform Capabilities at a Glance",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_AMB, Pt(2))

    features = [
        ("[Eye]", "Face Auto-Extraction",
         "Automatically isolates and centers facial regions using OpenCV Haar Cascades before inference, eliminating background noise.",
         C_ACCENT_BLUE),
        ("[Brain]", "EfficientNet-B3 CNN",
         "State-of-the-art binary classification neural network pre-trained on ImageNet, fine-tuned on deepfake forensic datasets.",
         C_ACCENT_CYAN),
        ("[Heatmap]", "Grad-CAM Heatmaps",
         "Generates transparent activation heatmaps highlighting exact pixel regions that triggered the deepfake classification.",
         C_ACCENT_ROSE),
        ("[Video]", "Video Frame Forensics",
         "Samples sequential frames, computes per-frame fake probability, and aggregates confidence ratios with temporal metrics.",
         C_ACCENT_EMR),
        ("[PDF]", "Forensic PDF Reports",
         "Compiles case metadata, media snapshots, Grad-CAM attention maps, and breakdown metrics into printable evidence documents.",
         C_ACCENT_AMB),
        ("[Chart]", "Analytics Dashboard",
         "Interactive Chart.js donut and bar visualizations tracking historical detection distributions and KPI counters in real-time.",
         C_ACCENT_CYAN),
        ("[CSV]", "CSV Data Export",
         "1-Click download of complete historical case registry in standard RFC 4180 CSV format for downstream analysis.",
         C_ACCENT_EMR),
        ("[API]", "Headless REST API",
         "POST /api/predict endpoint returns structured JSON results, enabling integration with mobile apps and microservices.",
         C_ACCENT_ROSE),
    ]

    cols = 4
    card_w = Inches(3.1)
    card_h = Inches(2.1)
    gap_x  = Inches(0.2)
    gap_y  = Inches(0.2)
    start_x = Inches(0.35)
    start_y = Inches(1.6)

    for i, (icon, title, desc, col) in enumerate(features):
        row = i // cols
        col_i = i % cols
        x = start_x + col_i * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h, C_BG_CARD)
        add_line(slide, x, y, card_w, col, Pt(4))
        add_text(slide, icon, x + Inches(0.12), y + Inches(0.1),
                 Inches(0.7), Inches(0.45), font_size=12, bold=True, color=col)
        add_text(slide, title, x + Inches(0.85), y + Inches(0.12),
                 card_w - Inches(1.0), Inches(0.4),
                 font_size=13, bold=True, color=col)
        add_text(slide, desc, x + Inches(0.12), y + Inches(0.6),
                 card_w - Inches(0.2), card_h - Inches(0.75),
                 font_size=11.5, color=C_TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MODEL PERFORMANCE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_performance(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_CYAN)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_CYAN)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_CYAN)
    add_text(slide, "07", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "MODEL PERFORMANCE", Inches(0.4), Inches(0.2), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_CYAN)
    add_text(slide, "Empirical Evaluation & Benchmark Metrics",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_CYAN, Pt(2))

    add_text(slide, "Evaluated on FaceForensics++ & Celeb-DF benchmarks",
             Inches(0.4), Inches(1.5), Inches(8), Inches(0.35),
             font_size=13, italic=True, color=C_TEXT_MUTED)

    # Big metric cards
    metrics = [
        ("69.86%", "Overall\nAccuracy",     C_ACCENT_EMR),
        ("74.49%", "Deepfake\nRecall",      C_ACCENT_CYAN),
        ("65.77%", "Authentic\nPrecision",  C_ACCENT_BLUE),
        ("0.684",  "Macro\nF1-Score",       C_ACCENT_AMB),
        ("0.584",  "Validation\nLoss",      C_ACCENT_ROSE),
    ]
    m_x = Inches(0.4)
    for val, label, col in metrics:
        add_rect(slide, m_x, Inches(1.95), Inches(2.4), Inches(1.6), C_BG_CARD)
        add_line(slide, m_x, Inches(1.95), Inches(2.4), col, Pt(4))
        add_text(slide, val, m_x, Inches(2.08), Inches(2.4), Inches(0.75),
                 font_size=32, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, label, m_x, Inches(2.8), Inches(2.4), Inches(0.65),
                 font_size=13, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)
        m_x += Inches(2.55)

    # Horizontal bar chart (simulated)
    add_rect(slide, Inches(0.4), Inches(3.75), Inches(6.2), Inches(3.45), C_BG_CARD)
    add_line(slide, Inches(0.4), Inches(3.75), Inches(6.2), C_ACCENT_CYAN, Pt(3))
    add_text(slide, "Performance Breakdown",
             Inches(0.55), Inches(3.85), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT_CYAN)

    bars = [
        ("Overall Accuracy",    69.86, C_ACCENT_EMR),
        ("Deepfake Recall",     74.49, C_ACCENT_CYAN),
        ("Authentic Precision", 65.77, C_ACCENT_BLUE),
        ("F1-Score (x100)",     68.4,  C_ACCENT_AMB),
    ]
    bar_y = Inches(4.35)
    bar_max_w = Inches(4.5)
    for label, pct, col in bars:
        add_text(slide, label, Inches(0.55), bar_y, Inches(1.9), Inches(0.3),
                 font_size=11, color=C_TEXT_MUTED)
        add_rect(slide, Inches(2.5), bar_y + Inches(0.04),
                 bar_max_w * (pct / 100), Inches(0.28), col)
        add_text(slide, f"{pct}%", Inches(2.5) + bar_max_w * (pct/100) + Inches(0.05),
                 bar_y, Inches(0.6), Inches(0.3),
                 font_size=11, bold=True, color=col)
        bar_y += Inches(0.62)

    # Confusion / Class metrics table
    add_rect(slide, Inches(6.8), Inches(3.75), Inches(6.15), Inches(3.45), C_BG_CARD)
    add_line(slide, Inches(6.8), Inches(3.75), Inches(6.15), C_ACCENT_AMB, Pt(3))
    add_text(slide, "Classification Report",
             Inches(6.95), Inches(3.85), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT_AMB)

    headers = ["Class", "Precision", "Recall", "F1"]
    col_xs  = [Inches(7.0), Inches(8.5), Inches(10.0), Inches(11.4)]
    h_y     = Inches(4.3)
    for hdr, cx in zip(headers, col_xs):
        add_text(slide, hdr, cx, h_y, Inches(1.3), Inches(0.3),
                 font_size=12, bold=True, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)

    rows = [
        ("REAL",      "65.77%", "63.28%", "0.645", C_ACCENT_EMR),
        ("FAKE",      "72.56%", "74.49%", "0.735", C_ACCENT_ROSE),
        ("Macro Avg", "69.17%", "68.88%", "0.690", C_TEXT_MUTED),
        ("Weighted",  "69.36%", "69.86%", "0.693", C_TEXT_MUTED),
    ]
    row_y = Inches(4.65)
    for cls, prec, rec, f1, col in rows:
        vals = [cls, prec, rec, f1]
        for v, cx in zip(vals, col_xs):
            add_text(slide, v, cx, row_y, Inches(1.3), Inches(0.38),
                     font_size=12, color=col, align=PP_ALIGN.CENTER)
        add_line(slide, Inches(7.0), row_y + Inches(0.38), Inches(5.8),
                 C_DIVIDER, Pt(1))
        row_y += Inches(0.42)

    add_text(slide, "Dataset: FaceForensics++ / Celeb-DF  |  Version: EfficientNet-B3 v3",
             Inches(0.4), Inches(7.2), Inches(12), Inches(0.28),
             font_size=10, italic=True, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WEB INTERFACE & REST API
# ═══════════════════════════════════════════════════════════════════════════════
def slide_interface(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_ROSE)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_ROSE)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_ROSE)
    add_text(slide, "08", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "WEB INTERFACE & REST API", Inches(0.4), Inches(0.2), Inches(8), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_ROSE)
    add_text(slide, "Forensic Web Dashboard & Programmatic Integration",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=28, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_ROSE, Pt(2))

    # Route table
    add_rect(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(5.6), C_BG_CARD)
    add_line(slide, Inches(0.4), Inches(1.55), Inches(6.0), C_ACCENT_ROSE, Pt(3))
    add_text(slide, "Web Routes",
             Inches(0.55), Inches(1.65), Inches(5), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT_ROSE)

    routes = [
        ("GET",  "/",             "Drag-and-Drop Forensic Ingestion Studio",  C_ACCENT_EMR),
        ("POST", "/predict",      "Processes image/video -> verdict + heatmap", C_ACCENT_CYAN),
        ("GET",  "/dashboard",    "KPI counters & Chart.js analytics",         C_ACCENT_BLUE),
        ("GET",  "/history",      "Searchable historical case registry",        C_ACCENT_AMB),
        ("GET",  "/export_csv",   "1-Click CSV case history download",          C_ACCENT_ROSE),
        ("GET",  "/architecture", "Interactive pipeline architecture view",     C_ACCENT_CYAN),
        ("GET",  "/results",      "Model benchmarks & confusion matrix",        C_ACCENT_EMR),
        ("GET",  "/about",        "Mission overview & tech specifications",     C_ACCENT_BLUE),
        ("GET",  "/health",       "System heartbeat health check",              C_ACCENT_AMB),
    ]
    r_y = Inches(2.15)
    for method, route, desc, col in routes:
        add_rect(slide, Inches(0.55), r_y, Inches(0.55), Inches(0.32), col)
        add_text(slide, method, Inches(0.55), r_y, Inches(0.55), Inches(0.32),
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, route, Inches(1.15), r_y, Inches(1.5), Inches(0.32),
                 font_size=11, bold=True, color=C_TEXT_MAIN)
        add_text(slide, desc, Inches(2.75), r_y, Inches(3.5), Inches(0.32),
                 font_size=10, color=C_TEXT_MUTED)
        r_y += Inches(0.49)

    # API Response JSON box
    add_rect(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(5.6), C_BG_CARD)
    add_line(slide, Inches(6.65), Inches(1.55), Inches(6.3), C_ACCENT_CYAN, Pt(3))
    add_text(slide, "REST API  --  POST /api/predict",
             Inches(6.8), Inches(1.65), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT_CYAN)

    add_text(slide, "Example JSON Response (Image):", Inches(6.8), Inches(2.1),
             Inches(5.8), Inches(0.3), font_size=12, bold=True, color=C_TEXT_MUTED)

    json_text = (
        '{\n'
        '  "success": true,\n'
        '  "media_type": "Image",\n'
        '  "filename": "suspect_face.jpg",\n'
        '  "prediction": "Fake",\n'
        '  "confidence": 89.64,\n'
        '  "fake_probability": 89.64,\n'
        '  "real_probability": 10.36,\n'
        '  "risk": "Very High",\n'
        '  "heatmap_url": "/uploads/hm.jpg",\n'
        '  "report_download_url": "/download/rpt.pdf"\n'
        '}'
    )
    add_rect(slide, Inches(6.8), Inches(2.45), Inches(5.95), Inches(2.55), C_DIVIDER)
    add_text(slide, json_text, Inches(6.9), Inches(2.5),
             Inches(5.7), Inches(2.5), font_size=10.5, color=C_ACCENT_CYAN)

    add_text(slide, "Example JSON Response (Video):", Inches(6.8), Inches(5.05),
             Inches(5.8), Inches(0.3), font_size=12, bold=True, color=C_TEXT_MUTED)
    json2 = (
        '{\n'
        '  "prediction": "Fake",  "confidence": 74.49,\n'
        '  "fake_frames": 26,     "real_frames": 9,\n'
        '  "fake_percentage": 74.29,  "risk": "High",\n'
        '  "report_download_url": "/download/report_vid.pdf"\n'
        '}'
    )
    add_rect(slide, Inches(6.8), Inches(5.4), Inches(5.95), Inches(1.35), C_DIVIDER)
    add_text(slide, json2, Inches(6.9), Inches(5.45),
             Inches(5.7), Inches(1.25), font_size=10.5, color=C_ACCENT_EMR)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DEPLOYMENT & INFRASTRUCTURE
# ═══════════════════════════════════════════════════════════════════════════════
def slide_deployment(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_BLUE)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_BLUE)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_BLUE)
    add_text(slide, "09", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "DEPLOYMENT & INFRASTRUCTURE", Inches(0.4), Inches(0.2), Inches(8), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_BLUE)
    add_text(slide, "Containerised & Cloud-Ready Deployment",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_BLUE, Pt(2))

    deploy_options = [
        ("Windows -- 1-Click", "run.bat",
         "Double-click to automatically verify Python dependencies,\nlaunch Flask server, and open browser.",
         C_ACCENT_EMR,  Inches(0.4),  Inches(1.6)),
        ("Linux / macOS",     "start.sh",
         "Bash shell script equivalent for Unix-based systems.\npython app.py -> http://127.0.0.1:5000",
         C_ACCENT_CYAN, Inches(4.6),  Inches(1.6)),
        ("Docker Container",  "docker-compose up",
         "Full containerised stack via Docker Compose.\nIsolated environment with automatic dependency management.",
         C_ACCENT_BLUE, Inches(8.8),  Inches(1.6)),
    ]

    for title, cmd, desc, col, x, y in deploy_options:
        add_rect(slide, x, y, Inches(3.95), Inches(2.2), C_BG_CARD)
        add_line(slide, x, y, Inches(3.95), col, Pt(4))
        add_text(slide, title, x + Inches(0.12), y + Inches(0.1),
                 Inches(3.7), Inches(0.4), font_size=14, bold=True, color=col)
        add_rect(slide, x + Inches(0.12), y + Inches(0.55),
                 Inches(3.7), Inches(0.38), C_DIVIDER)
        add_text(slide, f"$ {cmd}", x + Inches(0.2), y + Inches(0.57),
                 Inches(3.5), Inches(0.34), font_size=12, bold=True, color=C_ACCENT_AMB)
        add_text(slide, desc, x + Inches(0.12), y + Inches(1.05),
                 Inches(3.7), Inches(1.0), font_size=12, color=C_TEXT_MUTED)

    # Dockerfile highlights
    add_rect(slide, Inches(0.4), Inches(4.0), Inches(6.0), Inches(3.2), C_BG_CARD)
    add_line(slide, Inches(0.4), Inches(4.0), Inches(6.0), C_ACCENT_BLUE, Pt(3))
    add_text(slide, "Dockerfile Highlights",
             Inches(0.55), Inches(4.1), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT_BLUE)
    dockerfile_text = (
        "FROM python:3.11-slim\n"
        "WORKDIR /app\n"
        "COPY requirements.txt .\n"
        "RUN pip install --no-cache-dir -r requirements.txt\n"
        "COPY . .\n"
        "EXPOSE 5000\n"
        'CMD ["python", "app.py"]'
    )
    add_rect(slide, Inches(0.55), Inches(4.55), Inches(5.7), Inches(2.5), C_DIVIDER)
    add_text(slide, dockerfile_text, Inches(0.65), Inches(4.6),
             Inches(5.5), Inches(2.4), font_size=11, color=C_ACCENT_CYAN)

    # Cloud deployment
    add_rect(slide, Inches(6.65), Inches(4.0), Inches(6.3), Inches(3.2), C_BG_CARD)
    add_line(slide, Inches(6.65), Inches(4.0), Inches(6.3), C_ACCENT_AMB, Pt(3))
    add_text(slide, "render.yaml (Cloud Deployment)",
             Inches(6.8), Inches(4.1), Inches(5.8), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT_AMB)

    cloud_points = [
        "Platform: Render.com (render.yaml configured)",
        "Auto-builds from Git push -> zero-downtime deploy",
        "Gunicorn WSGI production server (not Flask dev)",
        "Environment: PYTHON_VERSION=3.11",
        "Persistent disk mount for SQLite case database",
        "Health check endpoint: GET /health",
    ]
    add_bullets(slide, cloud_points, Inches(6.8), Inches(4.55),
                Inches(6.0), Inches(2.5), font_size=12, color=C_TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
def slide_roadmap(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_EMR)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_EMR)
    add_rect(slide, Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35), C_ACCENT_EMR)
    add_text(slide, "10", Inches(12.5), Inches(0.15), Inches(0.7), Inches(0.35),
             font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "FUTURE ROADMAP", Inches(0.4), Inches(0.2), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=C_ACCENT_EMR)
    add_text(slide, "Planned Enhancements & Research Directions",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.75),
             font_size=30, bold=True, color=C_TEXT_MAIN)
    add_line(slide, Inches(0.4), Inches(1.4), Inches(8), C_ACCENT_EMR, Pt(2))

    phases = [
        ("Phase 1\nNear-Term", [
            "EfficientNet-V2 / ViT transformer backbone upgrade",
            "Multi-face detection in group photos & crowd scenes",
            "Audio-visual deepfake detection (voice + face sync)",
            "Batch upload API for bulk forensic processing",
        ], C_ACCENT_BLUE),
        ("Phase 2\nMid-Term", [
            "Federated learning across distributed investigation nodes",
            "Blockchain evidence timestamping (Ethereum / IPFS)",
            "Real-time video stream analysis (RTSP / WebRTC)",
            "Mobile companion app (React Native iOS + Android)",
        ], C_ACCENT_CYAN),
        ("Phase 3\nLong-Term", [
            "Cross-modal manipulation detection (text, audio, video)",
            "Zero-shot generalization to unseen GAN architectures",
            "Certified adversarial robustness against evasion attacks",
            "Integration with law enforcement forensic frameworks",
        ], C_ACCENT_AMB),
    ]

    p_x = Inches(0.4)
    for phase, items, col in phases:
        add_rect(slide, p_x, Inches(1.6), Inches(4.1), Inches(5.55), C_BG_CARD)
        add_line(slide, p_x, Inches(1.6), Inches(4.1), col, Pt(4))
        add_text(slide, phase, p_x + Inches(0.12), Inches(1.7),
                 Inches(3.85), Inches(0.65), font_size=15, bold=True, color=col)
        b_y = Inches(2.4)
        for item in items:
            add_rect(slide, p_x + Inches(0.12), b_y,
                     Inches(3.85), Inches(1.05), C_DIVIDER)
            add_line(slide, p_x + Inches(0.12), b_y, Inches(0.06), col, Pt(105))
            add_text(slide, item, p_x + Inches(0.3), b_y + Inches(0.05),
                     Inches(3.6), Inches(1.0), font_size=12, color=C_TEXT_MUTED)
            b_y += Inches(1.15)
        p_x += Inches(4.3)

    # Accuracy improvement note
    add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.55), Inches(0.3), C_BG_CARD)
    add_text(slide,
             "Target: >85% accuracy | >90% deepfake recall | <5% false-positive rate on production datasets",
             Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.28),
             font_size=12, color=C_ACCENT_CYAN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSION / THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    fill_bg(slide, C_BG_DARK)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT_CYAN)
    add_rect(slide, 0, 0, Inches(0.06), SLIDE_H, C_ACCENT_CYAN)
    add_rect(slide, SLIDE_W - Inches(0.06), 0, Inches(0.06), SLIDE_H, C_ACCENT_CYAN)
    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), C_ACCENT_CYAN)

    add_text(slide, "[ DIS ]",
             Inches(0), Inches(1.1), SLIDE_W, Inches(1.0),
             font_size=54, bold=True, color=C_ACCENT_CYAN, align=PP_ALIGN.CENTER)

    add_text(slide, "Thank You",
             Inches(0), Inches(2.2), SLIDE_W, Inches(0.95),
             font_size=48, bold=True, color=C_TEXT_MAIN, align=PP_ALIGN.CENTER)

    add_line(slide, Inches(3.5), Inches(3.15), Inches(6.3), C_ACCENT_CYAN, Pt(3))

    add_text(slide,
             "Deepfake Investigation System  --  AI-Powered Digital Forensics Platform",
             Inches(0), Inches(3.35), SLIDE_W, Inches(0.5),
             font_size=16, italic=True, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Summary bullets
    summary = [
        "EfficientNet-B3 deep CNN achieving 69.86% accuracy & 74.49% deepfake recall",
        "Explainable AI via Grad-CAM activation heatmaps for transparent verdicts",
        "Frame-by-frame video forensics with temporal anomaly metrics",
        "Auto-generated forensic PDF investigation reports",
        "Real-time analytics dashboard & REST API integration",
        "Containerised with Docker, deployable to Render.com cloud",
    ]
    s_y = Inches(3.95)
    for s in summary:
        add_text(slide, f"  >>  {s}", Inches(2.2), s_y, Inches(8.9), Inches(0.38),
                 font_size=14, color=C_TEXT_MUTED, align=PP_ALIGN.LEFT)
        s_y += Inches(0.43)

    add_text(slide, "Built with  TensorFlow  |  OpenCV  |  Flask  |  ReportLab  |  SQLAlchemy  |  Chart.js  |  Docker",
             Inches(0), Inches(7.1), SLIDE_W, Inches(0.32),
             font_size=11, color=C_TEXT_MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs);       print("  [OK] Slide 1  -- Title")
    slide_problem(prs);     print("  [OK] Slide 2  -- Problem Statement")
    slide_objectives(prs);  print("  [OK] Slide 3  -- Objectives")
    slide_architecture(prs);print("  [OK] Slide 4  -- System Architecture")
    slide_model(prs);       print("  [OK] Slide 5  -- Deep Learning Model")
    slide_features(prs);    print("  [OK] Slide 6  -- Key Features")
    slide_performance(prs); print("  [OK] Slide 7  -- Model Performance")
    slide_interface(prs);   print("  [OK] Slide 8  -- Web Interface & REST API")
    slide_deployment(prs);  print("  [OK] Slide 9  -- Deployment & Infrastructure")
    slide_roadmap(prs);     print("  [OK] Slide 10 -- Future Roadmap")
    slide_conclusion(prs);  print("  [OK] Slide 11 -- Conclusion / Thank You")

    prs.save(OUTPUT)
    print(f"\n[DONE] Presentation saved -> {OUTPUT}")


if __name__ == "__main__":
    main()
